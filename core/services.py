import os
import re
import json
import asyncio
from itertools import cycle
from io import BytesIO
from openai import AsyncOpenAI
from django.conf import settings

# Kutubxonalar
from docx import Document
from docx.shared import Pt, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Pt as PptxPt, Inches as PptxInches
from pptx.dml.color import RGBColor as PptxRGB
from pptx.enum.text import PP_ALIGN

# API KEYS
groq_keys = os.environ.get("GROQ_KEYS", "").split(",")
api_key_cycle = cycle([k for k in groq_keys if k])
GROQ_MODELS = ["llama-3.3-70b-versatile", "llama-3.1-70b-versatile"]

async def call_ai(prompt, json_mode=False):
    for _ in range(5):
        key = next(api_key_cycle)
        try:
            cl = AsyncOpenAI(api_key=key, base_url="https://api.groq.com/openai/v1")
            kwargs = {"model": GROQ_MODELS[0], "messages": [{"role": "user", "content": prompt}], "temperature": 0.7}
            if json_mode: kwargs["response_format"] = {"type": "json_object"}
            resp = await cl.chat.completions.create(**kwargs)
            await cl.close()
            return resp.choices[0].message.content
        except: continue
    return None

def clean_text(text):
    return re.sub(r'\*\*(.*?)\*\*', r'\1', text).replace("##", "").strip()

# --- DOCX ENGINE ---
def generate_docx(content_list, info):
    doc = Document()
    # Title
    for _ in range(4): doc.add_paragraph()
    p = doc.add_paragraph("O'ZBEKISTON RESPUBLIKASI OLIY TA'LIM VAZIRLIGI")
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER; p.runs[0].bold = True
    
    for _ in range(5): doc.add_paragraph()
    p = doc.add_paragraph(info['type'].upper()); p.alignment = WD_ALIGN_PARAGRAPH.CENTER; p.runs[0].font.size = Pt(24); p.runs[0].bold = True
    p = doc.add_paragraph(f"Mavzu: {info['topic']}"); p.alignment = WD_ALIGN_PARAGRAPH.CENTER; p.runs[0].font.size = Pt(16)
    
    for _ in range(6): doc.add_paragraph()
    p = doc.add_paragraph(f"Bajardi: {info['student']}"); p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    
    doc.add_page_break()
    
    for item in content_list:
        h = doc.add_paragraph(clean_text(item['title'])); h.alignment = WD_ALIGN_PARAGRAPH.CENTER; h.runs[0].bold = True; h.runs[0].font.size = Pt(16)
        p = doc.add_paragraph(clean_text(item['content'])); p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
    
    out = BytesIO(); doc.save(out); out.seek(0)
    return out

# --- PPTX ENGINE ---
def generate_pptx(content_list, info):
    prs = Presentation()
    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = info['topic']
    slide.placeholders[1].text = f"Tayyorladi: {info['student']}\nSana: {datetime.now().strftime('%d.%m.%Y')}"
    
    for item in content_list:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = clean_text(item['title'])
        tf = slide.placeholders[1].text_frame
        tf.text = clean_text(item['content'])
    
    out = BytesIO(); prs.save(out); out.seek(0)
    return out

# --- MAIN GENERATOR LOGIC ---
async def process_request(topic, pages, doc_type, student_name):
    # 1. Reja tuzish
    plan_prompt = f"Mavzu: {topic}. {pages} ta qismdan iborat reja tuz (JSON array formatida: [\"Sarlavha 1\", ...])."
    plan_json = await call_ai(plan_prompt, json_mode=True)
    try: titles = json.loads(plan_json).get('topics', []) if isinstance(json.loads(plan_json), dict) else json.loads(plan_json)
    except: titles = [f"{topic} - {i} qism" for i in range(1, pages+1)]
    
    if not isinstance(titles, list): titles = [f"{topic} haqida"]
    titles = titles[:pages]

    # 2. Tarkib yozish
    full_data = []
    for t in titles:
        if doc_type == 'pptx':
            prompt = f"Mavzu: {topic}. Slayd sarlavhasi: {t}. Slayd uchun 4-5 ta qisqa va lo'nda punkt yoz."
        else:
            prompt = f"Mavzu: {topic}. Bob: {t}. Shu bob uchun 500 so'zli ilmiy matn yoz."
        
        content = await call_ai(prompt) or "Ma'lumot topilmadi"
        full_data.append({"title": t, "content": content})
    
    # 3. Fayl yaratish
    info = {"topic": topic, "student": student_name, "type": "Taqdimot" if doc_type=='pptx' else "Referat"}
    if doc_type == 'pptx':
        return generate_pptx(full_data, info), f"{topic[:10]}.pptx"
    else:
        return generate_docx(full_data, info), f"{topic[:10]}.docx"