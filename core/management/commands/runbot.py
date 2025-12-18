import asyncio
import logging
import sys
import json
import re
import os
from io import BytesIO
from itertools import cycle
from datetime import datetime

from django.core.management.base import BaseCommand
from django.conf import settings
from asgiref.sync import sync_to_async

from aiogram import Bot, Dispatcher, F, types, Router
from aiogram.filters import CommandStart, Command, Filter
from aiogram.fsm.context import FSMContext
from aiogram.fsm.state import State, StatesGroup
from aiogram.fsm.storage.memory import MemoryStorage
from aiogram.types import (
    ReplyKeyboardMarkup, KeyboardButton, InlineKeyboardMarkup, InlineKeyboardButton,
    BufferedInputFile, CallbackQuery
)
from aiogram.utils.keyboard import InlineKeyboardBuilder

# --- DJANGO MODELLARNI IMPORT QILISH ---
from core.models import TGUser, DocHistory, Transaction

# --- CONFIG ---
from openai import AsyncOpenAI
BOT_TOKEN = os.environ.get("BOT_TOKEN")
ADMIN_ID = int(os.environ.get("ADMIN_ID", 0))
ADMIN_USERNAME = "admin"
KARTA_RAQAMI = os.environ.get("KARTA_RAQAMI", "8600 0000 0000 0000")
REFERRAL_BONUS = 10000

groq_keys_str = os.environ.get("GROQ_KEYS", "")
GROQ_API_KEYS = groq_keys_str.split(",") if "," in groq_keys_str else [groq_keys_str]
api_key_cycle = cycle([k for k in GROQ_API_KEYS if k])
GROQ_MODELS = ["llama-3.3-70b-versatile", "llama-3.1-70b-versatile"]

DEFAULT_PRICES = {
    "pptx_10": 5000, "pptx_15": 7000, "pptx_20": 10000,
    "docx_15": 5000, "docx_20": 7000, "docx_30": 12000
}

# --- KUTUBXONALAR ---
from docx import Document
from docx.shared import Pt, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Pt as PptxPt, Inches as PptxInches
from pptx.dml.color import RGBColor as PptxRGB
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from fpdf import FPDF

# FONT FIX
FONT_PATH = "DejaVuSans.ttf"
def check_font():
    if not os.path.exists(FONT_PATH):
        try:
            import requests
            url = "https://raw.githubusercontent.com/coreybutler/fonts/master/ttf/DejaVuSans.ttf"
            r = requests.get(url, timeout=30)
            with open(FONT_PATH, 'wb') as f: f.write(r.content)
        except: pass
check_font()

# ==============================================================================
# DJANGO BAZA BILAN ISHLASH (ASYNC WRAPPERS)
# ==============================================================================

@sync_to_async
def db_get_user(tg_id):
    try:
        return TGUser.objects.get(tg_id=tg_id)
    except TGUser.DoesNotExist:
        return None

@sync_to_async
def db_create_user(tg_id, username, full_name, referral_id=0):
    user, created = TGUser.objects.get_or_create(
        tg_id=tg_id,
        defaults={'username': username, 'full_name': full_name, 'referral_id': referral_id}
    )
    if not created:
        user.username = username
        user.full_name = full_name
        user.save()
    
    # Referal bonusi
    is_new = created
    if created and referral_id and referral_id != tg_id:
        try:
            referrer = TGUser.objects.get(tg_id=referral_id)
            referrer.balance += REFERRAL_BONUS
            referrer.invited_count += 1
            referrer.save()
            Transaction.objects.create(user=referrer, amount=REFERRAL_BONUS, type="referral_bonus")
        except: pass
    
    return is_new, user

@sync_to_async
def db_update_balance(tg_id, amount, type_trans="payment"):
    try:
        user = TGUser.objects.get(tg_id=tg_id)
        user.balance += amount
        user.save()
        Transaction.objects.create(user=user, amount=amount, type=type_trans)
        return True
    except: return False

@sync_to_async
def db_update_limit(tg_id, doc_type, amount):
    # doc_type: 'pptx' or 'docx' or 'pdf'
    try:
        user = TGUser.objects.get(tg_id=tg_id)
        if doc_type == 'pptx':
            user.free_pptx += amount
        elif doc_type == 'docx':
            user.free_docx += amount
        elif doc_type == 'pdf':
            user.free_pdf += amount
        user.save()
    except: pass

@sync_to_async
def db_save_history(tg_id, doc_type, topic, pages, info):
    try:
        user = TGUser.objects.get(tg_id=tg_id)
        # Modelga qarab maydonlar moslanadi
        # Agar core/models.py da student_name maydonlari bo'lmasa, faqat basic info saqlanadi
        DocHistory.objects.create(
            user=user,
            doc_type=doc_type,
            topic=topic,
            pages=pages
        )
    except: pass

@sync_to_async
def db_toggle_block(tg_id, status):
    try:
        user = TGUser.objects.get(tg_id=tg_id)
        user.is_blocked = status
        user.save()
        return True
    except: return False

@sync_to_async
def db_get_all_users():
    return list(TGUser.objects.values_list('tg_id', flat=True))

@sync_to_async
def db_get_admins():
    # Superuserlarni admin deb hisoblaymiz
    return list(TGUser.objects.filter(is_blocked=False).values_list('tg_id', flat=True)) 
    # Yoki agar ADMIN_ID env da bo'lsa:
    # return [ADMIN_ID]

@sync_to_async
def db_get_stats():
    total = TGUser.objects.count()
    blocked = TGUser.objects.filter(is_blocked=True).count()
    # Bugungi yangi userlar
    today = datetime.now().date()
    new_today = TGUser.objects.filter(joined_at__date=today).count()
    return total, blocked, new_today

# ==============================================================================
# ENGINES & AI (Oldingi koddan)
# ==============================================================================
def clean_text(text):
    text = re.sub(r'\*\*(.*?)\*\*', r'\1', text) 
    text = re.sub(r'##+', '', text)
    return re.sub(r'\n+', '\n', text).strip()

def extract_json_array(text):
    try:
        text = re.sub(r"```json", "", text).replace("```", "")
        start, end = text.find('['), text.rfind(']') + 1
        if start != -1 and end != -1: return json.loads(text[start:end])
        return []
    except: return []

PPTX_THEMES = {
    "modern_blue": {"bg": (240,248,255), "main": (0,51,102), "txt": (20,20,40), "shape": MSO_SHAPE.ROUNDED_RECTANGLE},
    "nature_green": {"bg": (240,255,240), "main": (34,139,34), "txt": (10,30,10), "shape": MSO_SHAPE.SNIP_2_DIAG_RECTANGLE},
    "creative_orange": {"bg": (255,250,240), "main": (255,69,0), "txt": (50,20,0), "shape": MSO_SHAPE.OVAL},
    "cyber_purple": {"bg": (20,0,30), "main": (0,255,255), "txt": (255,255,255), "shape": MSO_SHAPE.HEXAGON},
}

def create_presentation(data_list, info, design="modern_blue"):
    prs = Presentation()
    th = PPTX_THEMES.get(design, PPTX_THEMES["modern_blue"])
    bg_rgb, main_rgb, txt_rgb = PptxRGB(*th["bg"]), PptxRGB(*th["main"]), PptxRGB(*th["txt"])
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = bg_rgb
    shape = slide.shapes.add_shape(th['shape'], PptxInches(0.5), PptxInches(0.5), PptxInches(9), PptxInches(6.5))
    shape.fill.background(); shape.line.color.rgb = main_rgb; shape.line.width = PptxPt(4)
    tb = slide.shapes.add_textbox(PptxInches(1), PptxInches(2), PptxInches(8), PptxInches(2.5))
    p = tb.text_frame.add_paragraph(); p.text = info['topic'].upper(); p.font.size = PptxPt(40); p.font.bold = True; p.font.color.rgb = main_rgb; p.alignment = PP_ALIGN.CENTER
    ib = slide.shapes.add_textbox(PptxInches(1), PptxInches(5), PptxInches(8), PptxInches(2))
    p = ib.text_frame.add_paragraph(); p.text = f"Bajardi: {info['student']}\nGuruh: {info['group']}\nQabul qildi: {info['teacher']}"; p.font.size = PptxPt(18); p.font.color.rgb = txt_rgb; p.alignment = PP_ALIGN.CENTER

    for item in data_list:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid(); slide.background.fill.fore_color.rgb = bg_rgb
        head = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, PptxInches(10), PptxInches(1.2))
        head.fill.solid(); head.fill.fore_color.rgb = main_rgb
        ht = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(0.2), PptxInches(9), PptxInches(0.8))
        hp = ht.text_frame.add_paragraph(); hp.text = clean_text(item['title']); hp.font.size = PptxPt(32); hp.font.bold = True; hp.font.color.rgb = PptxRGB(255,255,255); hp.alignment = PP_ALIGN.CENTER
        bt = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(1.5), PptxInches(9), PptxInches(5.5))
        tf = bt.text_frame; tf.word_wrap = True
        content = clean_text(item['content'])
        fs = 24 if len(content) < 200 else 16
        for line in content.split('\n'):
            if len(line.strip()) > 2: p = tf.add_paragraph(); p.text = "• " + line.strip(); p.font.size = PptxPt(fs); p.font.color.rgb = txt_rgb; p.space_after = PptxPt(10)
    out = BytesIO(); prs.save(out); out.seek(0); return out

def create_document(data_list, info, doc_type="Referat"):
    doc = Document()
    for _ in range(4): doc.add_paragraph()
    p = doc.add_paragraph("O'ZBEKISTON RESPUBLIKASI\nOLIY TA'LIM VAZIRLIGI"); p.alignment = WD_ALIGN_PARAGRAPH.CENTER; p.runs[0].bold = True
    if info['edu_place'] != "-": p = doc.add_paragraph(info['edu_place'].upper()); p.alignment = WD_ALIGN_PARAGRAPH.CENTER; p.runs[0].bold = True
    for _ in range(6): doc.add_paragraph()
    p = doc.add_paragraph(doc_type.upper()); p.alignment = WD_ALIGN_PARAGRAPH.CENTER; p.runs[0].font.size = Pt(22); p.runs[0].bold = True
    p = doc.add_paragraph(f"Mavzu: {info['topic']}"); p.alignment = WD_ALIGN_PARAGRAPH.CENTER; p.runs[0].bold = True
    for _ in range(5): doc.add_paragraph()
    table = doc.add_table(rows=5, cols=2); table.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    def fill_row(idx, label, val):
        if val != "-": cell = table.rows[idx].cells[1]; p = cell.paragraphs[0]; r = p.add_run(f"{label}: {val}"); r.bold = True
    fill_row(0, "Bajardi", info['student']); fill_row(1, "Guruh", info['group']); fill_row(2, "Fakultet", info['direction']); fill_row(3, "Fan", info['subject']); fill_row(4, "Qabul qildi", info['teacher'])
    doc.add_page_break()
    for item in data_list:
        h = doc.add_paragraph(clean_text(item['title'])); h.alignment = WD_ALIGN_PARAGRAPH.CENTER; h.runs[0].bold = True; h.runs[0].font.size = Pt(16)
        for para in clean_text(item['content']).split('\n'):
            if len(para) > 5: p = doc.add_paragraph(para); p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY; p.paragraph_format.first_line_indent = Cm(1.27)
    out = BytesIO(); doc.save(out); out.seek(0); return out

class PDF(FPDF):
    def footer(self):
        self.set_y(-15); self.set_font("DejaVu", '', 10); self.cell(0, 10, f'{self.page_no()}', align='C')
def create_pdf(data_list, info, doc_type="Referat"):
    pdf = PDF()
    try: pdf.add_font("DejaVu", "", FONT_PATH, uni=True); pdf.add_font("DejaVu", "B", FONT_PATH, uni=True)
    except: return None
    pdf.set_font("DejaVu", "", 12); pdf.add_page()
    pdf.set_font("DejaVu", "B", 14); pdf.cell(0, 8, "O'ZBEKISTON RESPUBLIKASI", ln=True, align='C')
    pdf.ln(40); pdf.set_font("DejaVu", "B", 24); pdf.cell(0, 10, doc_type.upper(), ln=True, align='C')
    pdf.ln(10); pdf.set_font("DejaVu", "B", 16); pdf.multi_cell(0, 10, f"Mavzu: {info['topic']}", align='C')
    pdf.ln(40); pdf.set_font("DejaVu", "", 14); start_x = 100
    def add_line(label, val):
        if val != "-": pdf.set_x(start_x); pdf.set_font("DejaVu", "B", 14); pdf.cell(0, 10, f"{label}: {val}", ln=True)
    add_line("Bajardi", info['student']); add_line("Guruh", info['group']); add_line("Qabul qildi", info['teacher'])
    pdf.add_page()
    for item in data_list:
        pdf.set_font("DejaVu", "B", 16); pdf.multi_cell(0, 10, clean_text(item['title']), align='C'); pdf.ln(5)
        pdf.set_font("DejaVu", "", 12); pdf.multi_cell(0, 7, clean_text(item['content'])); pdf.ln(10)
    out = BytesIO(); out.write(pdf.output()); out.seek(0); return out

async def call_groq(messages):
    if not GROQ_API_KEYS: return None
    for _ in range(5):
        key = next(api_key_cycle)
        for model in GROQ_MODELS:
            try:
                cl = AsyncOpenAI(api_key=key, base_url="https://api.groq.com/openai/v1")
                resp = await cl.chat.completions.create(model=model, messages=messages, temperature=0.7, max_tokens=2500)
                await cl.close()
                return resp.choices[0].message.content
            except: continue
    return None

async def generate_full_content(topic, pages, doc_type, custom_plan, status_msg):
    async def progress(pct, text):
        try: await status_msg.edit_text(f"⏳ <b>Jarayon: {pct}%</b>\n\n⚙️ {text}", parse_mode="HTML")
        except: pass

    await progress(5, "Reja tuzilmoqda...")
    if doc_type == "taqdimot":
        prompt = f"Mavzu: {topic}. {pages} ta slayd uchun qiziqarli sarlavhalar (JSON array). Faqat JSON."
        res = await call_groq([{"role":"system","content":"JSON only."}, {"role":"user","content":prompt}])
        titles = extract_json_array(res)
        if not titles: titles = [f"{topic} - {i}-qism" for i in range(1, pages+1)]
        data = []
        for i, t in enumerate(titles[:pages]):
            await progress(10 + int((i/len(titles))*85), f"Slayd: {t}")
            p_text = f"Mavzu: {topic}. Slayd: {t}. 150-200 so'zdan iborat matn yoz."
            content = await call_groq([{"role":"user", "content":p_text}])
            data.append({"title": t, "content": content or "..."})
        return data
    else:
        num = max(6, int(pages/2) + 2)
        prompt = f"Mavzu: {topic}. {num} ta bobdan iborat reja."
        res = await call_groq([{"role":"user", "content":prompt}])
        chapters = [x.strip() for x in res.split('\n') if len(x)>5][:num]
        data = []
        for i, ch in enumerate(chapters):
            await progress(10 + int((i/len(chapters))*85), f"Bob: {ch}")
            p_text = f"Mavzu: {topic}. Bob: {ch}. 800 so'zli ilmiy matn yoz."
            content = await call_groq([{"role":"user", "content":p_text}])
            data.append({"title": ch, "content": content or "..."})
        return data

# ==============================================================================
# BOT HANDLERS (TO'LIQ)
# ==============================================================================
router = Router()

main_kb = ReplyKeyboardMarkup(keyboard=[[KeyboardButton(text="📊 Taqdimot"), KeyboardButton(text="📝 Mustaqil ish")], [KeyboardButton(text="📑 Referat"), KeyboardButton(text="💰 Balans & Referal")], [KeyboardButton(text="💳 To'lov qilish"), KeyboardButton(text="📞 Yordam")]], resize_keyboard=True)
cancel_kb = ReplyKeyboardMarkup(keyboard=[[KeyboardButton(text="❌ Bekor qilish")]], resize_keyboard=True)
skip_kb = InlineKeyboardMarkup(inline_keyboard=[[InlineKeyboardButton(text="➡️ O'tkazib yuborish", callback_data="skip")]])

class Form(StatesGroup):
    type = State(); topic = State(); plan = State(); student = State(); uni = State(); fac = State(); grp = State(); subj = State(); teach = State(); design = State(); len = State(); format = State()
class PayState(StatesGroup): screenshot = State(); amount = State()
class AdminState(StatesGroup): block_uid=State(); unblock_uid=State()
class IsAdmin(Filter): 
    async def __call__(self, m: types.Message): return await is_admin_check(m.from_user.id)

async def is_admin_check(user_id):
    return user_id == ADMIN_ID # Hozircha faqat bitta super admin

@router.message(F.text == "❌ Bekor qilish")
async def cancel_all(m: types.Message, state: FSMContext):
    await state.clear(); await m.answer("✅ Jarayon bekor qilindi.", reply_markup=main_kb)

@router.message(CommandStart())
async def start(m: types.Message, command: CommandObject):
    try:
        ref_id = int(command.args) if command.args and command.args.isdigit() else 0
    except: ref_id = 0
    
    is_new, user = await db_create_user(m.from_user.id, m.from_user.username, m.from_user.full_name, ref_id)
    
    if user.is_blocked: 
        await m.answer("🚫 Bloklangansiz.")
        return

    txt = "👋 <b>Assalomu alaykum!</b>\nMen professional akademik yordamchiman.\n\nWeb orqali foydalanish: https://edubot-pro.onrender.com"
    if is_new and ref_id: await m.bot.send_message(ref_id, f"🎉 Referal bonusi: {REFERRAL_BONUS} so'm!", parse_mode="HTML")
    await m.answer(txt, parse_mode="HTML", reply_markup=main_kb)

@router.message(F.text.in_(["📊 Taqdimot", "📝 Mustaqil ish", "📑 Referat"]))
async def start_order(m: types.Message, state: FSMContext):
    u = await db_get_user(m.from_user.id)
    if not u or u.is_blocked: return await m.answer("🚫 Bloklangansiz.")
    await state.update_data(dtype="taqdimot" if "Taqdimot" in m.text else "referat")
    await m.answer("📝 <b>Mavzuni yozing:</b>", parse_mode="HTML", reply_markup=cancel_kb); await state.set_state(Form.topic)

@router.message(Form.topic)
async def get_topic(m: types.Message, state: FSMContext):
    await state.update_data(topic=m.text); await m.answer("📋 Reja bormi?", reply_markup=skip_kb); await state.set_state(Form.plan)

@router.callback_query(F.data == "skip", Form.plan)
async def skip_p(c: CallbackQuery, state: FSMContext): await state.update_data(plan="-"); await c.message.answer("👤 Ism-Familiya:"); await state.set_state(Form.student)
@router.message(Form.plan)
async def get_plan(m: types.Message, state: FSMContext): await state.update_data(plan=m.text); await m.answer("👤 Ism-Familiya:"); await state.set_state(Form.student)
@router.message(Form.student)
async def get_s(m: types.Message, state: FSMContext): await state.update_data(student=m.text); await m.answer("🏫 Universitet:", reply_markup=skip_kb); await state.set_state(Form.uni)
@router.callback_query(F.data == "skip", Form.uni)
async def skip_u(c: CallbackQuery, state: FSMContext): await state.update_data(uni="-"); await c.message.answer("📚 Fakultet:", reply_markup=skip_kb); await state.set_state(Form.fac)
@router.message(Form.uni)
async def get_u(m: types.Message, state: FSMContext): await state.update_data(uni=m.text); await m.answer("📚 Fakultet:", reply_markup=skip_kb); await state.set_state(Form.fac)
@router.callback_query(F.data == "skip", Form.fac)
async def skip_f(c: CallbackQuery, state: FSMContext): await state.update_data(fac="-"); await c.message.answer("🔢 Guruh:", reply_markup=skip_kb); await state.set_state(Form.grp)
@router.message(Form.fac)
async def get_f(m: types.Message, state: FSMContext): await state.update_data(fac=m.text); await m.answer("🔢 Guruh:", reply_markup=skip_kb); await state.set_state(Form.grp)
@router.callback_query(F.data == "skip", Form.grp)
async def skip_g(c: CallbackQuery, state: FSMContext): await state.update_data(grp="-"); await c.message.answer("📘 Fan nomi:"); await state.set_state(Form.subj)
@router.message(Form.grp)
async def get_g(m: types.Message, state: FSMContext): await state.update_data(grp=m.text); await m.answer("📘 Fan nomi:"); await state.set_state(Form.subj)
@router.message(Form.subj)
async def get_sub(m: types.Message, state: FSMContext): await state.update_data(subj=m.text); await m.answer("👨‍🏫 O'qituvchi:"); await state.set_state(Form.teach)

@router.message(Form.teach)
async def get_teach(m: types.Message, state: FSMContext):
    await state.update_data(teacher=m.text)
    d = await state.get_data()
    kb = InlineKeyboardBuilder()
    if d['dtype'] == "taqdimot":
        [kb.button(text=th.title(), callback_data=f"d_{th}") for th in PPTX_THEMES.keys()]
        kb.adjust(2); await m.answer("🎨 Dizayn:", reply_markup=kb.as_markup()); await state.set_state(Form.design)
    else:
        await state.update_data(design="simple")
        kb.button(text="Word (.docx)", callback_data="fmt_docx"); kb.button(text="PDF (.pdf)", callback_data="fmt_pdf")
        await m.answer("📂 Format:", reply_markup=kb.as_markup()); await state.set_state(Form.format)

@router.callback_query(F.data.startswith("d_"), Form.design)
async def sel_d(c: CallbackQuery, state: FSMContext):
    await state.update_data(design=c.data[2:], fmt="pptx")
    kb = InlineKeyboardBuilder()
    for i in [10, 15, 20]: kb.button(text=f"{i} slayd", callback_data=f"len_{i}_{DEFAULT_PRICES[f'pptx_{i}']}")
    kb.row(InlineKeyboardButton(text="🔙 Orqaga", callback_data="back"))
    await c.message.edit_text("📄 Hajm:", reply_markup=kb.as_markup()); await state.set_state(Form.len)

@router.callback_query(F.data.startswith("fmt_"), Form.format)
async def sel_f(c: CallbackQuery, state: FSMContext):
    await state.update_data(fmt=c.data[4:])
    kb = InlineKeyboardBuilder()
    for i in [15, 20, 30]: kb.button(text=f"{i} bet", callback_data=f"len_{i}_{DEFAULT_PRICES[f'docx_{i}']}")
    kb.row(InlineKeyboardButton(text="🔙 Orqaga", callback_data="back"))
    await c.message.edit_text("📄 Hajm:", reply_markup=kb.as_markup()); await state.set_state(Form.len)

@router.callback_query(F.data == "back", Form.len)
async def back_h(c: CallbackQuery, state: FSMContext):
    d = await state.get_data(); await c.answer()
    if d['dtype'] == "taqdimot": 
        kb = InlineKeyboardBuilder(); [kb.button(text=th.title(), callback_data=f"d_{th}") for th in PPTX_THEMES.keys()]
        await c.message.edit_text("🎨 Dizayn:", reply_markup=kb.adjust(2).as_markup()); await state.set_state(Form.design)
    else: 
        kb = InlineKeyboardBuilder(); kb.button(text="Word", callback_data="fmt_docx"); kb.button(text="PDF", callback_data="fmt_pdf")
        await c.message.edit_text("📂 Format:", reply_markup=kb.as_markup()); await state.set_state(Form.format)

@router.callback_query(F.data.startswith("len_"), Form.len)
async def generate(c: CallbackQuery, state: FSMContext):
    await c.message.delete()
    _, p_s, c_s = c.data.split("_"); pages, cost = int(p_s), int(c_s)
    uid = c.from_user.id; u = await db_get_user(uid); d = await state.get_data()
    fmt = d.get('fmt', 'pptx')
    
    limit_key = 'free_pptx' if fmt == 'pptx' else ('free_docx' if fmt == 'docx' else 'free_pdf')
    current_limit = getattr(u, limit_key, 0)
    is_free = current_limit > 0
    
    if not is_free and u.balance < cost: return await c.message.answer(f"❌ Mablag' yetarli emas! {cost} so'm kerak.", reply_markup=main_kb)
    
    msg = await c.message.answer("⏳ Tayyorlanmoqda...")
    content = await generate_full_content(d['topic'], pages, d['dtype'], d['plan'], msg)
    if not content: return await msg.edit_text("❌ Xatolik.")
    
    info = {k: d.get(k, "-") for k in ['topic','student','uni','fac','grp','subj','teacher']}; info['edu_place']=d.get('uni','-'); info['direction']=d.get('fac','-'); info['group']=d.get('grp','-'); info['subject']=d.get('subj','-')
    
    if fmt == "pptx": f = create_presentation(content, info, d.get('design','modern_blue')); fn = f"{d['topic'][:20]}.pptx"
    elif fmt == "pdf": f = create_pdf(content, info, d['dtype']); fn = f"{d['topic'][:20]}.pdf"
    else: f = create_document(content, info, d['dtype']); fn = f"{d['topic'][:20]}.docx"
    
    await c.message.answer_document(BufferedInputFile(f.read(), filename=fn), caption="✅ Tayyor!", reply_markup=main_kb)
    await msg.delete()
    
    if is_free: await db_update_limit(uid, fmt, -1) # NOTE: db_update_limit logic needs to match field names
    else: await db_update_balance(uid, -cost, "service_fee")
    await db_save_history(uid, d['dtype'], d['topic'], pages, info)
    await state.clear()

# --- PAYMENT ---
@router.message(F.text == "💳 To'lov qilish")
async def pay(m: types.Message):
    kb = InlineKeyboardBuilder(); [kb.button(text=f"💎 {x}", callback_data=f"pay_{x}") for x in [5000, 10000, 15000]]; kb.adjust(2)
    kb.row(InlineKeyboardButton(text="❌ Yopish", callback_data="close"))
    await m.answer("👇 Summani tanlang:", reply_markup=kb.as_markup())

@router.callback_query(F.data.startswith("pay_"))
async def pay_i(c: CallbackQuery, state: FSMContext):
    await state.update_data(amount=int(c.data.split("_")[1])); await state.set_state(PayState.screenshot)
    await c.message.edit_text(f"💳 Karta: <code>{KARTA_RAQAMI}</code>\n📸 Chekni yuboring.")

@router.message(PayState.screenshot, F.photo)
async def pay_c(m: types.Message, state: FSMContext):
    amt = (await state.get_data())['amount']; kb = InlineKeyboardBuilder()
    kb.button(text="✅", callback_data=f"ap_{m.from_user.id}_{amt}"); kb.button(text="❌", callback_data=f"de_{m.from_user.id}")
    try: await m.bot.send_photo(ADMIN_ID, m.photo[-1].file_id, caption=f"To'lov: {amt}\nUser: {m.from_user.full_name}", reply_markup=kb.as_markup())
    except: pass
    await m.answer("✅ Yuborildi.", reply_markup=main_kb); await state.clear()

@router.callback_query(F.data.startswith("ap_"))
async def ap(c: CallbackQuery):
    if c.from_user.id != ADMIN_ID: return
    _, u, a = c.data.split("_"); await db_update_balance(int(u), int(a)); await c.message.delete(); await c.bot.send_message(int(u), "✅ To'lov qabul qilindi.")

@router.callback_query(F.data.startswith("de_"))
async def de(c: CallbackQuery):
    if c.from_user.id != ADMIN_ID: return
    await c.message.delete(); await c.bot.send_message(int(c.data.split("_")[1]), "❌ To'lov rad etildi.")

@router.callback_query(F.data == "close")
async def cl(c: CallbackQuery): await c.message.delete()

@router.message(F.text == "💰 Balans & Referal")
async def bal(m: types.Message):
    u = await db_get_user(m.from_user.id)
    await m.answer(f"💰 Balans: {u.balance}\n👥 Takliflar: {u.invited_count}\n🎁 Limit: PPTX-{u.free_pptx} DOCX-{u.free_docx}", parse_mode="HTML")

# --- ADMIN PANEL ---
def admin_kb_markup():
    kb = InlineKeyboardBuilder()
    kb.button(text="📊 Statistika", callback_data="adm_stats")
    kb.button(text="🚫 Bloklash", callback_data="adm_block")
    kb.button(text="✅ Ochish", callback_data="adm_unblock")
    kb.button(text="🚪 Yopish", callback_data="close")
    return kb.as_markup()

@router.message(Command("admin"), IsAdmin())
async def adm_cmd(m: types.Message):
    await m.answer("Admin Panel:", reply_markup=admin_kb_markup())

@router.callback_query(F.data == "adm_stats", IsAdmin())
async def ast(c: CallbackQuery): 
    t, b, n = await db_get_stats()
    await c.message.edit_text(f"📊 Jami: {t}\n🚫 Blok: {b}\n🆕 Bugun: {n}", reply_markup=admin_kb_markup())

@router.callback_query(F.data == "adm_block", IsAdmin())
async def b_ui(c: CallbackQuery, state: FSMContext):
    await c.message.answer("User ID:"); await state.set_state(AdminState.block_uid)

@router.message(AdminState.block_uid, IsAdmin())
async def b_do(m: types.Message, state: FSMContext):
    if await db_toggle_block(int(m.text), True): await m.answer("🚫 Bloklandi.")
    await state.clear()

@router.callback_query(F.data == "adm_unblock", IsAdmin())
async def ub_ui(c: CallbackQuery, state: FSMContext):
    await c.message.answer("User ID:"); await state.set_state(AdminState.unblock_uid)

@router.message(AdminState.unblock_uid, IsAdmin())
async def ub_do(m: types.Message, state: FSMContext):
    if await db_toggle_block(int(m.text), False): await m.answer("✅ Ochildi.")
    await state.clear()

class Command(BaseCommand):
    def handle(self, *args, **options):
        bot = Bot(token=BOT_TOKEN)
        dp = Dispatcher(storage=MemoryStorage())
        dp.include_router(router)
        asyncio.run(dp.start_polling(bot))